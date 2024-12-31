from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from imagegen import generate_image
from pptx.shapes.picture import Picture


class PresentationAPI:
    def __init__(self, file_path='presentation.pptx'):
        """Initialize the PresentationAPI with the specified file path.
        
        Args:
            file_path (str): The path to the PowerPoint file. If it doesn't exist, a new presentation will be created.
        """
        self.file_path = file_path
        try:
            self.prs = Presentation(file_path)
        except Exception as e:
            self.prs = Presentation()
            print(f"Error loading presentation: {e}")
        
        self.current_slide = len(self.prs.slides) - 1 if len(self.prs.slides) > 0 else 0  

    def create_slide(self, layout='title', title='{}', content='{}', subtitle=None):
        """Create a new slide in the presentation and update the current slide index.
        
        Args:
            layout (str): The layout type for the slide.
            title (str): The title of the slide.
            content (str): The content of the slide.
            subtitle (str, optional): The subtitle of the slide.
        
        Returns:
            int: The index of the newly created slide.
        """
        layout_dict = {
            'title': 0, 'title and content': 1, 'section header': 2,
            'two content': 3, 'comparison': 4, 'title only': 5,
            'blank': 6, 'content with caption': 7, 'picture with caption': 8
        }
        slide_layout = self.prs.slide_layouts[layout_dict[layout.lower()]]
        slide = self.prs.slides.add_slide(slide_layout)

        if title != '{}':
            title_placeholder = slide.placeholders[0]
            title_placeholder.text = title

        if subtitle and layout.lower() == 'title':
            subtitle_placeholder = slide.placeholders[1]
            subtitle_placeholder.text = subtitle

        if content != '{}' and layout.lower() != 'title':
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = content

        self.save_presentation()
        self.current_slide = len(self.prs.slides) -1  # Update current slide index
        return self.current_slide

    def next_slide(self):
        """Move to the next slide in the presentation."""
        if self.current_slide < len(self.prs.slides)-1:
            self.current_slide += 1
        return self.current_slide

    def previous_slide(self):
        """Move to the previous slide in the presentation."""
        if self.current_slide > 0:
            self.current_slide -= 1
        return self.current_slide

    def get_current_slide(self):
        """Get the index of the current slide."""
        return self.current_slide

    def add_text(self, text, title_placeholder_num=0, content_placeholder_num=1):
        """Add non-bullet text to the current slide.

        Args:
            text (str): The text to add.
            title_placeholder_num (int): The index of the title placeholder (default is 0).
            content_placeholder_num (int): The index of the content placeholder (default is 1).
        """
        slide = self.prs.slides[self.current_slide]  # Use the current slide index

        if title_placeholder_num < len(slide.placeholders):
            title_placeholder = slide.placeholders[title_placeholder_num]
            top_position = title_placeholder.top + title_placeholder.height
        else:
            # Posición alternativa si el placeholder de título no está disponible
            top_position = Pt(100)  # Puedes ajustar esta posición según tus necesidades

        # Intentar obtener el placeholder de contenido
        if content_placeholder_num < len(slide.placeholders):
            content_placeholder = slide.placeholders[content_placeholder_num]
            left = content_placeholder.left
            width = content_placeholder.width
            height = content_placeholder.height
            
        else:
            # Posición y tamaño alternativos si el placeholder de contenido no está disponible
            left = Pt(100)
            width = Pt(400)
            height = Pt(300)
        
        textbox = slide.shapes.add_textbox(left, top_position, width, height)  # Ajusta la altura según tus necesidades
        text_frame = textbox.text_frame
        text_frame.word_wrap=True

        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(32)
        p.alignment = PP_ALIGN.CENTER

        self.save_presentation()
    
    def add_bullet_points(self, slide_num, points, placeholder=1):
        """Add bullet points to a slide.

        Args:
            slide_num (int): The slide number to add bullet points to.
            points (list): A list of strings representing the bullet points.
            placeholder (int): The index of the placeholder for bullet points (default is 1).
        """
        if slide_num > len(self.prs.slides) or slide_num < 1:
            print(f"Error: slide {slide_num} does not exist.")
            return
        slide = self.prs.slides[slide_num]
        text_frame = slide.placeholders[placeholder].text_frame
        text_frame.clear()  # Clear previous content

        for point in points:
            p = text_frame.add_paragraph()  # Add a new paragraph
            p.text = point
            p.level = 0  # Set bullet level
        self.save_presentation()

    def edit_text(self, slide_num, placeholder, new_text, paragraph_index=0):#creates a * to the left
        """Edit existing text in a slide's placeholder.

        Args:
            slide_num (int): The slide number to edit.
            placeholder (int): The index of the placeholder to edit.
            new_text (str): The new text to replace the existing text.
            paragraph_index (int): The index of the paragraph to edit (default is 0).
        """
        slide = self.prs.slides[slide_num]
        text_frame = slide.placeholders[placeholder].text_frame
        if paragraph_index < len(text_frame.paragraphs):
            text_frame.paragraphs[paragraph_index].text = new_text
            self.save_presentation()
        else:
            print(f"Paragraph index {paragraph_index} does not exist in slide {slide_num}, placeholder {placeholder}")

    def delete_text(self, slide_num, placeholder, paragraph_index=None):
        """Delete text from a slide's placeholder.

        Args:
            slide_num (int): The slide number to delete text from.
            placeholder (int): The index of the placeholder to delete text from.
            paragraph_index (int, optional): The index of the paragraph to delete. If None, all text will be deleted.
        """
        slide = self.prs.slides[slide_num]
        text_frame = slide.placeholders[placeholder].text_frame
        if paragraph_index is None:
            text_frame.clear()
        elif paragraph_index < len(text_frame.paragraphs):
            text_frame._element.remove(text_frame.paragraphs[paragraph_index]._element)
        else:
            print(f"Paragraph index {paragraph_index} does not exist in slide {slide_num}, placeholder {placeholder}")
        self.save_presentation()

    def change_font(self, slide_num, placeholder_num, font_name, font_size):
        """Change the font of text in a slide's placeholder.

        Args:
            slide_num (int): The slide number to change the font.
            placeholder_num (int): The index of the placeholder to change the font.
            font_name (str): The name of the font to apply.
            font_size (int): The size of the font to apply.
        """
        slide = self.prs.slides[slide_num]
        text_frame = slide.placeholders[placeholder_num].text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = Pt(font_size)
        self.save_presentation()

    def change_background_color(self, slide_num, color):
        """Change the background color of a slide.

        Args:
            slide_num (int): The slide number to change the background color.
            color (tuple): A tuple of RGB values (R, G, B) for the background color.
        """
        slide = self.prs.slides[slide_num]
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color)
        self.save_presentation()

    def save_presentation(self):
        """Save the presentation to the specified file path."""
        self.prs.save(self.file_path)
        print(f"Presentation saved as '{self.file_path}'")

    def get_slide_count(self):
        """Get the total number of slides in the presentation.

        Returns:
            int: The total number of slides.
        """
        return len(self.prs.slides)

    def add_image(self, slide_num, prompt, image_path='generated_image.png'):
        """Add an image to a slide.

        Args:
            slide_num (int): The slide number to add the image to.
            prompt (str): The prompt for generating the image.
            image_path (str): The path where the generated image is saved.
        """
        # Assuming generate_image saves the image to image_path
        image_generated = generate_image(prompt)

        # Check if the image was generated successfully
        if image_generated:
            # Add the image to the slide
            slide = self.prs.slides[slide_num]

            # Position and size adjustments for the image
            left = Inches(5)
            top = Inches(2.5)
            width = Inches(4)
            height = Inches(3)

            # Insert the image into the slide
            slide.shapes.add_picture(image_path, left, top, width, height)

            # Save the presentation
            self.save_presentation()

    def delete_images(self, slide_num):
        """Delete all images from a slide.

        Args:
            slide_num (int): The slide number to delete images from.
        """
        slide = self.prs.slides[slide_num]
        # Create a list of shapes that are pictures
        images_to_delete = [shape for shape in slide.shapes if isinstance(shape, Picture)]

        # Remove the images from the slide
        for image in images_to_delete:
            sp = image._element
            sp.getparent().remove(sp)

        self.save_presentation()

    def add_picture_with_caption(self, slide_num, prompt, title, text, image_path='generated_image.png'):
        """
        Add a generated picture with a caption to a slide based on a prompt.

        Args:
            slide_num (int): The slide number where the picture and caption will be added.
            prompt (str): The prompt to generate the image.
            caption_text (str): The caption text to add below the image.
        """
        if slide_num >= len(self.prs.slides) or slide_num < 0:
            print(f"Error: slide {slide_num} does not exist.")
            return
        
        image_generated=generate_image(prompt)

        if image_generated:
        # Add the image to the slide
            slide = self.prs.slides[slide_num]

            image_placeholder = slide.shapes[1]
            image_placeholder.fill.solid()
            left = image_placeholder.left
            top = image_placeholder.top
            width = image_placeholder.width
            height = image_placeholder.height
            slide.shapes.add_picture(image_path, left, top, width, height)

            title_placeholder = slide.shapes.title
            title_placeholder.text = title
            caption_placeholder = slide.placeholders[2]  
            caption_placeholder.text = text

        self.save_presentation()


    def delete_slide(self, slide_num):
        """Delete a slide from the presentation.
        
        Args:
            slide_num (int): The index of the slide to delete.
        """
        xml_slides = self.prs.slides._sldIdLst  # Lista de IDs de diapositivas
        slides = list(xml_slides)  # Convertimos a lista para manipular

        if 0 <= slide_num < len(slides):
            xml_slides.remove(slides[slide_num])  # Eliminamos la diapositiva
            self.prs.save(self.file_path)
        
                
        self.save_presentation()
        print(f"Slide {slide_num} deleted successfully.")



    def save_presentation(self):
        """Save the presentation to the specified file path."""
        self.prs.save(self.file_path)
        print(f"Presentation saved as '{self.file_path}'")

    # Add methods for other functionalities (add_bullet_points, edit_text, etc.)



# Example usage
api = PresentationAPI('presentation.pptx')



# ayudame a crear una presentacion unica pero totalmente diferent  como el siguiente ejemplo para anadirlo a este dataset



# {

#     "instruction": "Create a title slide",

#     "input": {

#       "title": "",

#       "content": ""

#     },

#     "output": "### CODE\napi.create_slide(layout='title')\n### END_CODE"

#   },

#   {

#     "instruction": "Change the title to 'Sports'",

#     "input": {

#       "title": "Sports",

#       "content": ""

#     },

#     "output": "### CODE\napi.edit_text(api.get_current_slide(), placeholder=0, new_text='Sports')\n### END_CODE"

#   },

#   {

#     "instruction": "Add a title and content slide",

#     "input": {

#       "title": "",

#       "content": ""

#     },

#     "output": "### CODE\napi.create_slide(layout='title and content')\n### END_CODE"

#   },

#   {

#     "instruction": "Change the title to 'Advantages of Sports'",

#     "input": {

#       "title": "Advantages of Sports",

#       "content": ""

#     },

#     "output": "### CODE\napi.edit_text(api.get_current_slide(), placeholder=0, new_text='Advantages of Sports')\n### END_CODE"

#   },

#   {

#     "instruction": "Add content related to the title",

#     "input": {

#       "title": "",

#       "content": "Sports Improve physical health and fitness, Develop teamwork and social skills, Reduce stress and anxiety, Promote discipline and dedication, Improve cognitive function"

#     },

#     "output": "### CODE\napi.add_text('Sports Improve physical health and fitness, Develop teamwork and social skills, Reduce stress and anxiety, Promote discipline and dedication, Improve cognitive function', 0, 1)\n### END_CODE"

#   },

#   {

#     "instruction": "Create another title and content slide with the title 'Types of Sports'",

#     "input": {

#       "title": "Types of Sports",

#       "content": ""

#     },

#     "output": "### CODE\napi.create_slide(layout='title and content', title='Types of Sports')\n### END_CODE"

#   },

#   {

#     "instruction": "Create five bullet points with different types of sports",

#     "input": {

#       "title": "",

#       "content": [

#         "Soccer - The world's most popular sport",

#         "Basketball - A fast-paced indoor/outdoor sport",

#         "Tennis - An individual or doubles racket sport",

#         "Swimming - A full-body workout sport",

#         "Athletics - Track and field events"

#       ]

#     },

#     "output": "### CODE\napi.add_bullet_points(api.get_current_slide(), points=['Soccer - The world\\'s most popular sport', 'Basketball - A fast-paced indoor/outdoor sport', 'Tennis - An individual or doubles racket sport', 'Swimming - A full-body workout sport', 'Athletics - Track and field events'])\n### END_CODE"

#   },

#   {

#     "instruction": "Create a title and content slide with the title 'Football' and include an image of a football player",

#     "input": {

#       "title": "Football",

#       "content": "Football is a team sport played between two teams of eleven players. It is the world's most popular sport, followed by billions of people worldwide."

#     },

#     "output": "### CODE\napi.create_slide(layout='picture with caption')\napi.add_picture_with_caption(api.get_current_slide(), prompt='professional soccer player in action', title='Football', text='Football is a team sport played between two teams of eleven players. It is the world\\'s most popular sport, followed by billions of people worldwide.')\n### END_CODE"

#   },

#   {

#     "instruction": "Create a title and content slide with the title 'Volleyball' and include content that discusses volleyball",

#     "input": {

#       "title": "Volleyball",

#       "content": "Is a Dynamic team sport played with six players per team, requires excellent teamwork and coordination"

#     },

#     "output": "### CODE\napi.create_slide(layout='title and content', title='Volleyball', content='Is a Dynamic team sport played with six players per team, requires excellent teamwork and coordination')\n### END_CODE"

#   },

#   {

#     "instruction": "Add a volleyball image",

#     "input": {

#       "title": "",

#       "content": "a volleyball team playing"

#     },

#     "output": "### CODE\napi.add_image(api.get_current_slide(), 'a volleyball team playing')\n### END_CODE"

#   },

#   {

#     "instruction": "Create a slide with the title 'Video Games,' include a text that discusses video games, and finally add a related image",

#     "input": {

#       "title": "Video Games",

#       "content": "While not traditional sports, competitive gaming and esports have emerged as a significant part of modern sporting culture, requiring quick reflexes, strategic thinking, and team coordination."

#     },

#     "output": "### CODE\napi.create_slide(layout='picture with caption')\napi.add_picture_with_caption(api.get_current_slide(), prompt='esports gaming setup with professional gamer', title='Video Games', text='While not traditional sports, competitive gaming and esports have emerged as a significant part of modern sporting culture, requiring quick reflexes, strategic thinking, and team coordination.')\n### END_CODE"

#   },

#   {

#     "instruction": "Delete a slide that I didn't like",

#     "input": {

#       "title": "",

#       "content": ""

#     },

#     "output": "### CODE\napi.delete_slide(api.get_current_slide())\n### END_CODE"

#   },

#   {

#     "instruction": "Create a slide with the title 'Conclusions' that provides a conclusion about sports",

#     "input": {

#       "title": "Conclusions",

#       "content": [

#         "Sports are essential for physical and mental well-being",

#         "They provide valuable life lessons and skills",

#         "Regular participation leads to better health outcomes",

#         "Sports bring people together across cultures",

#         "Both traditional sports and esports have their place in modern society"

#       ]

#     },

#     "output": "### CODE\napi.create_slide(layout='title and content', title='Conclusions')\napi.add_bullet_points(api.get_current_slide(), points=['Sports are essential for physical and mental well-being', 'They provide valuable life lessons and skills', 'Regular participation leads to better health outcomes', 'Sports bring people together across cultures', 'Both traditional sports and esports have their place in modern society'])\n### END_CODE"

#   },

# {

#     "instruction": "Change the title of slide 3 to 'Nutrition Basics'",

#     "input": {

#       "title": "Nutrition Basics",

#       "content": ""

#     },

#     "output": "### CODE\napi.edit_text(3-1, placeholder=0, new_text='Nutrition Basics')\n### END_CODE"

#   },

#   {

#     "instruction": "Update the content of slide 2 with new bullet points",

#     "input": {

#       "title": "",

#       "content": [

#         "Eat more vegetables",

#         "Reduce processed foods",

#         "Stay hydrated",

#         "Control portions"

#       ]

#     },

#     "output": "### CODE\napi.add_bullet_points(2-1, points=['Eat more vegetables', 'Reduce processed foods', 'Stay hydrated', 'Control portions'])\n### END_CODE"

#   },



# # Create title slide
# api.create_slide(layout='title')

# # Change the title to "Sports"
# api.edit_text(api.get_current_slide(), placeholder=0, new_text="Sports")

# # Add a title and content slide
# api.create_slide(layout='title and content')

# # Change the title to "Advantages of Sports"
# api.edit_text(api.get_current_slide(), placeholder=0, new_text="Advantages of Sports")

# # Add content related to the title
# api.add_text("Sports Improve physical health and fitness, Develop teamwork and social skills, Reduce stress and anxiety, Promote discipline and dedication, Improve cognitive function",0,1,)

# # Create another title and content slide with the title "types of sports"
# api.create_slide(layout='title and content', title="Types of Sports")

# # Create five bullet points with different types of sports
# api.add_bullet_points(api.get_current_slide(), points=["Soccer - The world's most popular sport", "Basketball - A fast-paced indoor/outdoor sport", "Tennis - An individual or doubles racket sport", "Swimming - A full-body workout sport", "Athletics - Track and field events"])

# # Create a title and content slide with the title "Football" and include an image of a football player
# api.create_slide(layout='picture with caption')
# api.add_picture_with_caption(
#     api.get_current_slide(),
#     prompt="professional soccer player in action",
#     title="Football",
#     text="Football is a team sport played between two teams of eleven players. It is the world's most popular sport, followed by billions of people worldwide."
# )

# # Create a title and content slide with the title "Volleyball" and include content that discusses volleyball
# api.create_slide(layout='title and content', title="Volleyball", content='Is a Dynamic team sport played with six players per team, requires excellent teamwork and coordination')

# # add the a volleyball image
# api.add_image(api.get_current_slide(), 'a volleyball team playing')

# # Create a slide with the title "Video Games," include a text that discusses video games, and finally add a related image

# api.create_slide(layout='picture with caption')
# api.add_picture_with_caption(
#     api.get_current_slide(),
#     prompt="esports gaming setup with professional gamer",
#     title="Video Games",
#     text="While not traditional sports, competitive gaming and esports have emerged as a significant part of modern sporting culture, requiring quick reflexes, strategic thinking, and team coordination."
# )

# # Delete a slide that I didn't like

# api.delete_slide(api.get_current_slide())

# # Create a slide with the title "Conclusions" that provides a conclusion about sports

# api.create_slide(layout='title and content', title="Conclusions")
# api.add_bullet_points(api.get_current_slide(), points=[
#     "Sports are essential for physical and mental well-being",
#     "They provide valuable life lessons and skills",
#     "Regular participation leads to better health outcomes",
#     "Sports bring people together across cultures",
#     "Both traditional sports and esports have their place in modern society"
# ])



# api.create_slide('picture with caption')
# api.add_picture_with_caption(api.get_current_slide(), 'giagiant forest', 'Forest2' ,'Descriptivo para la imagen.')
# api.edit_text(api.get_current_slide(), placeholder=1, new_text="Contenido de la columna derecha")
# code_block = "api.create_slide(layout='title', title='Technologies')\napi.add_text(api.get_current_slide(), 'Exploring the latest advancements.')"

# # Ejecutar el código
# exec_locals = {}  # Diccionario local para guardar variables creadas por el código
# exec(code_block, globals(), exec_locals)
# # api.save_presentation()
# api.create_slide('title','hello')
# current=api.get_current_slide()
# print(current)
#   # Move to the next slide
# api.next_slide()
# api.next_slide()
# api.next_slide()
# api.add_bullet_points(7,['halloooo333o'])

